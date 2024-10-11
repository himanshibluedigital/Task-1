#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Fri Oct 11 10:59:10 2024

@author: himanshikumar
"""

import pytesseract
from PIL import Image
import os
from datetime import datetime
import fitz
import PyPDF2
import docx
import pptx
import csv
import pandas as pd

def extract_text_from_image(image_path):
    # Extract text from the image using Tesseract-OCR
    text = pytesseract.image_to_string(Image.open(image_path))
    return text

def extract_text_from_pdf(pdf_path):
    # Extract text from the PDF using PyMuPDF
    with fitz.open(pdf_path) as doc:
        text = ''
        for page in doc:
            text += page.get_text()
        return text

def extract_text_from_docx(docx_path):
    # Extract text from the DOCX using python-docx
    doc = docx.Document(docx_path)
    text = ''
    for para in doc.paragraphs:
        text += para.text
    return text

def extract_text_from_pptx(pptx_path):
    # Extract text from the PPTX using python-pptx
    presentation = pptx.Presentation(pptx_path)
    text = ''
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text
    return text

def extract_text_from_csv(csv_path):
    # Extract text from the CSV using pandas
    df = pd.read_csv(csv_path)
    text = ''
    for column in df.columns:
        text += column + '\n'
        for value in df[column]:
            text += str(value) + '\n'
    return text

def extract_text_from_document(document_path):
    # Determine the file type and extract text accordingly
    file_extension = document_path.split('.')[-1].lower()
    if file_extension == 'pdf':
        return extract_text_from_pdf(document_path)
    elif file_extension == 'docx':
        return extract_text_from_docx(document_path)
    elif file_extension == 'pptx':
        return extract_text_from_pptx(document_path)
    elif file_extension == 'csv':
        return extract_text_from_csv(document_path)
    elif file_extension == 'jpg' or file_extension == 'png':
        return extract_text_from_image(document_path)
    else:
        return 'Unsupported file type.'

# Test the function
document_path = '/Users/himanshikumar/Desktop/Exams.pdf'  # replace with your document path
text = extract_text_from_document(document_path)
print(text)